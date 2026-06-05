import io
import os
import re
import shutil
import subprocess
import tempfile
import zipfile
from pathlib import Path
from datetime import date, datetime
from numbers import Number
import math

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.ticker import FuncFormatter

import fitz  # PyMuPDF
import streamlit as st
from openpyxl import load_workbook
from openpyxl.utils import get_column_letter
from openpyxl.utils.datetime import from_excel
from openpyxl.worksheet.page import PageMargins
from PIL import Image, ImageChops, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


APP_TITLE = "Excel to PowerPoint Dashboard Builder"
DEFAULT_LOGO = Path("assets/default_logo.png")

BYU_NAVY = RGBColor(16, 49, 101)
SIDEBAR_BG = RGBColor(245, 246, 248)
SIDEBAR_DIVIDER = RGBColor(205, 208, 214)


def find_soffice() -> str | None:
    """Find the LibreOffice executable."""
    candidates = [
        shutil.which("soffice"),
        shutil.which("libreoffice"),
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    ]
    for candidate in candidates:
        if candidate and Path(candidate).exists():
            return str(candidate)
    return None


def save_upload(uploaded_file, destination: Path) -> Path:
    destination.parent.mkdir(parents=True, exist_ok=True)
    destination.write_bytes(uploaded_file.getbuffer())
    return destination


def sanitize_filename(name: str) -> str:
    allowed = "abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ0123456789-_ .()"
    return "".join(ch if ch in allowed else "_" for ch in name).strip() or "file"


CELL_RANGE_RE = re.compile(
    r"(?P<sheet>'(?:[^']|'')+'|[^!]+)!"
    r"\$?(?P<start_col>[A-Z]{1,3})\$?(?P<start_row>\d+)"
    r"(?::\$?(?P<end_col>[A-Z]{1,3})\$?(?P<end_row>\d+))?"
)


def parse_excel_date(value) -> date | None:
    """Return a Python date when a cell contains an Excel date or date-looking text."""
    if isinstance(value, datetime):
        return value.date()
    if isinstance(value, date):
        return value

    # Some workbooks store dates as Excel serial numbers. Keep the bounds
    # conservative so normal dashboard counts are not accidentally treated as dates.
    if isinstance(value, Number) and 25000 <= float(value) <= 90000:
        try:
            converted = from_excel(value)
            return converted.date() if isinstance(converted, datetime) else converted
        except Exception:
            pass

    if isinstance(value, str):
        text = value.strip()
        for fmt in ("%m/%d/%Y", "%m/%d/%y", "%Y-%m-%d", "%m-%d-%Y", "%Y/%m/%d"):
            try:
                return datetime.strptime(text, fmt).date()
            except ValueError:
                pass
    return None

def unquote_sheet_name(name: str) -> str:
    if name.startswith("'") and name.endswith("'"):
        return name[1:-1].replace("''", "'")
    return name


def quote_sheet_name(name: str) -> str:
    return "'" + name.replace("'", "''") + "'" if any(ch in name for ch in " !'()") else name


def get_horizontal_range_parts(formula: str):
    """Parse formulas like NittyGrittySheet!$B$23:$AE$23.

    Returns None unless the reference is a single horizontal row range.
    """
    if not formula:
        return None
    match = CELL_RANGE_RE.search(formula)
    if not match:
        return None

    start_row = int(match.group("start_row"))
    end_row = int(match.group("end_row") or match.group("start_row"))
    if start_row != end_row:
        return None

    from openpyxl.utils import column_index_from_string

    sheet_name = unquote_sheet_name(match.group("sheet"))
    start_col = column_index_from_string(match.group("start_col"))
    end_col = column_index_from_string(match.group("end_col") or match.group("start_col"))
    return sheet_name, start_row, start_col, end_col


def make_horizontal_range_formula(sheet_name: str, row: int, start_col: int, end_col: int) -> str:
    from openpyxl.utils import get_column_letter

    quoted = quote_sheet_name(sheet_name)
    return f"{quoted}!${get_column_letter(start_col)}${row}:${get_column_letter(end_col)}${row}"


def clear_chart_reference_cache(ref_obj) -> None:
    """Force Excel/LibreOffice to rebuild chart caches from the updated source range."""
    if not ref_obj:
        return
    for attr in ("numCache", "strCache", "multiLvlStrCache"):
        if hasattr(ref_obj, attr):
            try:
                setattr(ref_obj, attr, None)
            except Exception:
                pass


def latest_date_col_for_row(
    ws,
    row: int,
    start_col: int,
    end_col_hint: int | None = None,
    cap_at_today: bool = False,
) -> tuple[int | None, date | None]:
    """Find the newest date column in a horizontal date row.

    By default this uses the latest date that exists in the uploaded workbook,
    even if that date is later than the date when the app code was written. This
    is what makes the app work for future weekly updates like 5/22, 5/29, 6/5,
    and so on.

    If cap_at_today=True, it ignores dates after the server's current date. Use
    that only when the workbook has blank placeholder columns for future weeks.
    """
    max_col = max(ws.max_column, end_col_hint or 0)
    dated_cols: list[tuple[date, int]] = []

    for col in range(start_col, max_col + 1):
        cell_date = parse_excel_date(ws.cell(row=row, column=col).value)
        if not cell_date:
            continue
        if cap_at_today and cell_date > date.today():
            continue
        dated_cols.append((cell_date, col))

    if not dated_cols:
        return None, None

    latest_date, latest_col = max(dated_cols, key=lambda item: item[0])
    return latest_col, latest_date

def latest_used_col_for_row(ws, row: int, start_col: int, end_col_hint: int | None = None) -> int | None:
    max_col = max(ws.max_column, end_col_hint or 0)
    latest = None
    for col in range(start_col, max_col + 1):
        value = ws.cell(row=row, column=col).value
        if value not in (None, ""):
            latest = col
    return latest


def update_ref_formula_to_end_col(ref_obj, target_end_col: int | None) -> bool:
    if not ref_obj or not getattr(ref_obj, "f", None) or not target_end_col:
        return False

    parts = get_horizontal_range_parts(ref_obj.f)
    if not parts:
        return False

    sheet_name, row, start_col, current_end_col = parts
    if target_end_col <= current_end_col:
        # Still clear stale caches, because cached chart data can make LibreOffice
        # render old points even when the workbook data is newer.
        clear_chart_reference_cache(ref_obj)
        return False

    ref_obj.f = make_horizontal_range_formula(sheet_name, row, start_col, target_end_col)
    clear_chart_reference_cache(ref_obj)
    return True


def extend_charts_to_latest_available_date(wb, cap_at_today: bool = False) -> tuple[int, date | None]:
    """Extend horizontal chart ranges to the newest date in the uploaded workbook.

    This fixes the exact problem where the Excel file contains a newer weekly
    pull date, but the saved chart range still ends at an older date. The app now
    scans the chart's category/date row and extends the category and value ranges
    to the newest date it finds. Nothing is hardcoded to 5/15, 5/22, or any other
    specific date.
    """
    updates = 0
    newest_date_seen: date | None = None

    # Ask Excel/LibreOffice to rebuild formulas and chart caches during render.
    try:
        wb.calculation.calcMode = "auto"
        wb.calculation.fullCalcOnLoad = True
        wb.calculation.forceFullCalc = True
    except Exception:
        pass

    for chart_ws in wb.worksheets:
        for chart in getattr(chart_ws, "_charts", []):
            for series in getattr(chart, "series", []):
                target_end_col = None

                # Prefer the category/date row because it controls the weekly x-axis.
                cat = getattr(series, "cat", None)
                cat_refs = []
                if cat is not None:
                    cat_refs.extend([getattr(cat, "strRef", None), getattr(cat, "numRef", None)])

                for cat_ref in cat_refs:
                    if not cat_ref or not getattr(cat_ref, "f", None):
                        continue
                    parts = get_horizontal_range_parts(cat_ref.f)
                    if not parts:
                        continue
                    sheet_name, row, start_col, end_col = parts
                    if sheet_name not in wb.sheetnames:
                        continue
                    target_end_col, latest_seen = latest_date_col_for_row(
                        wb[sheet_name],
                        row,
                        start_col,
                        end_col_hint=end_col,
                        cap_at_today=cap_at_today,
                    )
                    if latest_seen and (newest_date_seen is None or latest_seen > newest_date_seen):
                        newest_date_seen = latest_seen
                    if target_end_col:
                        updates += int(update_ref_formula_to_end_col(cat_ref, target_end_col))
                        break

                # Extend each value series to the same final column as its date row.
                val = getattr(series, "val", None)
                val_ref = getattr(val, "numRef", None) if val is not None else None
                if val_ref and getattr(val_ref, "f", None):
                    if not target_end_col:
                        parts = get_horizontal_range_parts(val_ref.f)
                        if parts:
                            sheet_name, row, start_col, end_col = parts
                            if sheet_name in wb.sheetnames:
                                target_end_col = latest_used_col_for_row(
                                    wb[sheet_name], row, start_col, end_col_hint=end_col
                                )
                    updates += int(update_ref_formula_to_end_col(val_ref, target_end_col))

                # Clear title caches too, just in case LibreOffice uses stale chart XML.
                tx = getattr(series, "tx", None)
                tx_ref = getattr(tx, "strRef", None) if tx is not None else None
                clear_chart_reference_cache(tx_ref)

    normalize_weekly_chart_axes(wb)
    return updates, newest_date_seen

def prepare_workbook_for_rendering(
    source_xlsx: Path,
    output_xlsx: Path,
    include_sheets: list[str],
    fit_each_sheet_to_one_page: bool,
    margins: float,
    auto_extend_latest_date: bool,
    cap_chart_dates_at_today: bool,
) -> list[str]:
    """Copy workbook, hide excluded sheets, and set print layout."""
    wb = load_workbook(source_xlsx)

    visible_sheets = []
    for ws in wb.worksheets:
        if ws.title in include_sheets:
            ws.sheet_state = "visible"
            visible_sheets.append(ws.title)
        else:
            ws.sheet_state = "hidden"

    if not visible_sheets:
        raise ValueError("No sheets selected. Select at least one sheet.")

    # Excel requires at least one visible sheet; this is guaranteed above.
    for ws in wb.worksheets:
        if ws.title not in visible_sheets:
            continue

        min_r = min_c = 10**9
        max_r = max_c = 0
        for row in ws.iter_rows():
            for cell in row:
                if cell.value is not None:
                    min_r = min(min_r, cell.row)
                    min_c = min(min_c, cell.column)
                    max_r = max(max_r, cell.row)
                    max_c = max(max_c, cell.column)

        if max_r:
            ws.print_area = f"{get_column_letter(min_c)}{min_r}:{get_column_letter(max_c)}{max_r}"

        ws.page_setup.orientation = "landscape"
        ws.page_margins = PageMargins(
            left=margins,
            right=margins,
            top=margins,
            bottom=margins,
            header=0.1,
            footer=0.1,
        )

        if fit_each_sheet_to_one_page:
            ws.page_setup.fitToWidth = 1
            ws.page_setup.fitToHeight = 1
            ws.sheet_properties.pageSetUpPr.fitToPage = True

    if auto_extend_latest_date:
        extend_charts_to_latest_available_date(wb, cap_at_today=cap_chart_dates_at_today)

    wb.save(output_xlsx)
    return visible_sheets


def convert_to_pdf(input_file: Path, output_dir: Path, soffice_path: str) -> Path:
    """Convert a spreadsheet or PowerPoint to PDF using LibreOffice."""
    output_dir.mkdir(parents=True, exist_ok=True)
    cmd = [
        soffice_path,
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(output_dir),
        str(input_file),
    ]
    result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, timeout=180)
    if result.returncode != 0:
        raise RuntimeError(
            "LibreOffice failed to convert the file to PDF.\n\n"
            f"STDOUT:\n{result.stdout}\n\nSTDERR:\n{result.stderr}"
        )

    pdf_path = output_dir / f"{input_file.stem}.pdf"
    if not pdf_path.exists():
        pdfs = list(output_dir.glob("*.pdf"))
        if not pdfs:
            raise FileNotFoundError("LibreOffice did not create a PDF.")
        pdf_path = pdfs[0]
    return pdf_path


def pdf_to_images(pdf_path: Path, output_dir: Path, dpi: int) -> list[Path]:
    """Render every PDF page to a PNG image."""
    output_dir.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(str(pdf_path))
    image_paths = []
    zoom = dpi / 72
    matrix = fitz.Matrix(zoom, zoom)

    for page_index in range(len(doc)):
        page = doc[page_index]
        pix = page.get_pixmap(matrix=matrix, alpha=False)
        out_path = output_dir / f"page_{page_index + 1:02d}.png"
        pix.save(str(out_path))
        image_paths.append(out_path)
    doc.close()
    return image_paths


def crop_white_space(image_path: Path, output_path: Path, threshold: int = 12, margin_px: int = 18) -> Path:
    """Crop white margins around an image while keeping a small border."""
    image = Image.open(image_path).convert("RGB")
    white = Image.new("RGB", image.size, (255, 255, 255))
    diff = ImageChops.difference(image, white)
    mask = diff.point(lambda p: 255 if p > threshold else 0)
    bbox = mask.getbbox()

    if bbox is None:
        image.save(output_path)
        return output_path

    left = max(0, bbox[0] - margin_px)
    top = max(0, bbox[1] - margin_px)
    right = min(image.width, bbox[2] + margin_px)
    bottom = min(image.height, bbox[3] + margin_px)
    cropped = image.crop((left, top, right, bottom))
    cropped.save(output_path)
    return output_path


def resize_for_ppt(image_path: Path, output_path: Path, max_dimension: int = 2200) -> Path:
    """Downsize very large screenshots so the PPTX stays reasonably small."""
    image = Image.open(image_path).convert("RGB")
    scale = min(max_dimension / image.width, max_dimension / image.height, 1.0)
    if scale < 1.0:
        new_size = (int(image.width * scale), int(image.height * scale))
        image = image.resize(new_size, Image.LANCZOS)
    image.save(output_path, quality=92, optimize=True)
    return output_path


def make_circular_logo_transparent(source: Path, output: Path) -> Path:
    """Make the area outside the main dark circular logo transparent.

    This works well for the BYU Marriott circular logo. It keeps the white letters
    inside the circle because it uses an ellipse mask instead of removing all white pixels.
    """
    image = Image.open(source).convert("RGBA")
    rgb = image.convert("RGB")
    pixels = rgb.load()

    dark_points = []
    step = max(1, min(image.size) // 500)
    for y in range(0, image.height, step):
        for x in range(0, image.width, step):
            r, g, b = pixels[x, y]
            if b > r + 15 and b > g + 5 and (r + g + b) < 450:
                dark_points.append((x, y))

    if not dark_points:
        image.save(output)
        return output

    xs = [p[0] for p in dark_points]
    ys = [p[1] for p in dark_points]
    bbox = (min(xs), min(ys), max(xs), max(ys))
    pad = int(min(image.size) * 0.02)
    bbox = (
        max(0, bbox[0] - pad),
        max(0, bbox[1] - pad),
        min(image.width - 1, bbox[2] + pad),
        min(image.height - 1, bbox[3] + pad),
    )

    mask = Image.new("L", image.size, 0)
    draw = ImageDraw.Draw(mask)
    draw.ellipse(bbox, fill=255)
    image.putalpha(mask)
    image.save(output)
    return output


def add_generated_intro_slide(prs: Presentation, title: str, subtitle: str, logo_path: Path | None) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.75))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BYU_NAVY
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.85), Inches(1.65), Inches(8.2), Inches(1.0))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = BYU_NAVY

    subtitle_box = slide.shapes.add_textbox(Inches(0.85), Inches(2.65), Inches(7.8), Inches(1.0))
    p2 = subtitle_box.text_frame.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(45, 45, 45)

    if logo_path and logo_path.exists():
        slide.shapes.add_picture(str(logo_path), Inches(9.75), Inches(1.25), width=Inches(2.25), height=Inches(2.25))


def add_intro_from_template(prs: Presentation, template_pptx: Path, working_dir: Path, soffice_path: str) -> None:
    """Use the first slide from a template deck as a full-slide screenshot."""
    pdf_dir = working_dir / "template_pdf"
    pdf_path = convert_to_pdf(template_pptx, pdf_dir, soffice_path)
    page_images = pdf_to_images(pdf_path, working_dir / "template_pages", dpi=160)
    if not page_images:
        raise RuntimeError("Template PowerPoint did not render any slides.")

    first_slide_img = page_images[0]
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(str(first_slide_img), 0, 0, width=prs.slide_width, height=prs.slide_height)


def build_powerpoint(
    sheet_image_paths: list[Path],
    sheet_names: list[str],
    output_pptx: Path,
    logo_path: Path | None,
    intro_template_pptx: Path | None,
    working_dir: Path,
    soffice_path: str,
    show_sheet_name: bool,
    sidebar_width_inches: float,
    title: str,
) -> Path:
    prs = Presentation()
    prs.slide_width = 12192000  # 13.333 in
    prs.slide_height = 6858000  # 7.5 in
    blank = prs.slide_layouts[6]

    if intro_template_pptx:
        add_intro_from_template(prs, intro_template_pptx, working_dir, soffice_path)
    else:
        add_generated_intro_slide(
            prs,
            title=title,
            subtitle="Generated from the uploaded Excel dashboard.",
            logo_path=logo_path,
        )

    slide_w = 13.333
    slide_h = 7.5
    left_pad = 0.10
    top_pad = 0.27
    right_pad_before_sidebar = 0.15
    content_w = slide_w - sidebar_width_inches - left_pad - right_pad_before_sidebar
    content_h = 6.95

    for sheet_name, img_path in zip(sheet_names, sheet_image_paths):
        slide = prs.slides.add_slide(blank)

        sidebar_left = Inches(slide_w - sidebar_width_inches)
        sidebar_w = Inches(sidebar_width_inches)

        sidebar = slide.shapes.add_shape(1, sidebar_left, 0, sidebar_w, prs.slide_height)
        sidebar.fill.solid()
        sidebar.fill.fore_color.rgb = SIDEBAR_BG
        sidebar.line.fill.background()

        divider = slide.shapes.add_shape(1, sidebar_left, Inches(0.2), 1, Inches(7.0))
        divider.fill.solid()
        divider.fill.fore_color.rgb = SIDEBAR_DIVIDER
        divider.line.fill.background()

        if logo_path and logo_path.exists():
            logo_size = min(sidebar_width_inches - 0.22, 1.05)
            logo_left = slide_w - sidebar_width_inches + (sidebar_width_inches - logo_size) / 2
            slide.shapes.add_picture(
                str(logo_path),
                Inches(logo_left),
                Inches(0.28),
                width=Inches(logo_size),
                height=Inches(logo_size),
            )

        if show_sheet_name:
            box = slide.shapes.add_textbox(
                Inches(slide_w - sidebar_width_inches + 0.04),
                Inches(1.35),
                Inches(sidebar_width_inches - 0.08),
                Inches(4.8),
            )
            p = box.text_frame.paragraphs[0]
            p.text = "Nitty Gritty" if sheet_name == "NittyGrittySheet" else sheet_name
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = BYU_NAVY

        image = Image.open(img_path)
        iw, ih = image.size
        scale = min((content_w * 96) / iw, (content_h * 96) / ih)
        pic_w = Inches(iw / 96 * scale)
        pic_h = Inches(ih / 96 * scale)
        pic_left = Inches(left_pad) + (Inches(content_w) - pic_w) / 2
        pic_top = Inches(top_pad) + (Inches(content_h) - pic_h) / 2

        slide.shapes.add_picture(str(img_path), pic_left, pic_top, width=pic_w, height=pic_h)

    prs.save(output_pptx)
    return output_pptx


def make_zip(file_path: Path, zip_path: Path) -> Path:
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as archive:
        archive.write(file_path, arcname=file_path.name)
    return zip_path




def set_weekly_chart_axis_to_show_latest_label(chart) -> None:
    """Ask Excel/LibreOffice to draw all weekly category tick labels.

    LibreOffice sometimes extends the chart data but still auto-skips the newest
    x-axis label. This makes the final weekly pull, such as 5/22/2026, appear
    on the axis instead of looking like the old chart stopped at 5/15/2026.
    """
    try:
        chart.x_axis.tickLblSkip = 1
        chart.x_axis.tickMarkSkip = 1
        chart.x_axis.noMultiLvlLbl = False
    except Exception:
        pass


def is_weekly_horizontal_chart(chart) -> bool:
    for series in getattr(chart, "series", []):
        cat = getattr(series, "cat", None)
        if cat is None:
            continue
        for ref in (getattr(cat, "strRef", None), getattr(cat, "numRef", None)):
            if ref and getattr(ref, "f", None) and get_horizontal_range_parts(ref.f):
                return True
    return False


def normalize_weekly_chart_axes(wb) -> None:
    for ws in wb.worksheets:
        for chart in getattr(ws, "_charts", []):
            if is_weekly_horizontal_chart(chart):
                set_weekly_chart_axis_to_show_latest_label(chart)


def parse_numeric(value) -> float | None:
    if isinstance(value, Number):
        return float(value)
    if isinstance(value, str):
        text = value.strip().replace(",", "")
        if not text or text in {"-", "--", "—"}:
            return None
        is_percent = text.endswith("%")
        if is_percent:
            text = text[:-1].strip()
        try:
            parsed = float(text)
            return parsed / 100 if is_percent else parsed
        except ValueError:
            return None
    return None


def find_weekly_block_rows(wb, sheet_name: str) -> tuple[int, int, int, dict[str, int]] | None:
    """Return date/value rows for this dashboard's weekly charts.

    The placement dashboards store weekly chart data in NittyGrittySheet.
    For MSB, the block begins under "2026 MSB Weekly Placement". For each
    major, the block begins at a row whose first cell matches the worksheet name.
    """
    if "NittyGrittySheet" not in wb.sheetnames or sheet_name == "NittyGrittySheet":
        return None

    ws = wb["NittyGrittySheet"]

    if sheet_name == "MSB Full-Time Dashboard":
        for row in range(1, ws.max_row + 1):
            value = ws.cell(row=row, column=1).value
            if isinstance(value, str) and "MSB Weekly Placement" in value:
                return (
                    row + 1,
                    row + 2,
                    row + 4,
                    {
                        "Accepted an offer": row + 5,
                        "Actively seeking": row + 6,
                        "Not Reported": row + 8,
                    },
                )
        # Fallback for the known BYU Marriott dashboard layout.
        return (23, 24, 26, {"Accepted an offer": 27, "Actively seeking": 28, "Not Reported": 30})

    for row in range(1, ws.max_row + 1):
        value = ws.cell(row=row, column=1).value
        if isinstance(value, str) and value.strip() == sheet_name:
            return (
                row + 1,
                row + 2,
                row + 4,
                {
                    "Accepted an offer": row + 5,
                    "Actively seeking": row + 6,
                    "Not Reported": row + 8,
                },
            )

    return None


def extract_horizontal_weekly_series(
    ws,
    date_row: int,
    value_row: int,
    fallback_date_row: int | None = None,
) -> tuple[list[date], list[float]]:
    """Extract a date/value series across columns B:last.

    If a status date row was not copied forward but the value row has newer
    values, fallback_date_row lets the app borrow the matching date labels from
    the placement date row. This prevents stale chart screenshots when future
    weekly columns exist but one date row was not extended correctly.
    """
    dates: list[date] = []
    values: list[float] = []
    max_col = ws.max_column

    for col in range(2, max_col + 1):
        current_date = parse_excel_date(ws.cell(row=date_row, column=col).value)
        if current_date is None and fallback_date_row is not None:
            current_date = parse_excel_date(ws.cell(row=fallback_date_row, column=col).value)
        if current_date is None:
            continue

        current_value = parse_numeric(ws.cell(row=value_row, column=col).value)
        if current_value is None:
            continue

        dates.append(current_date)
        values.append(current_value)

    # Sort and dedupe by date, keeping the last value for duplicate dates.
    deduped: dict[date, float] = {}
    for d, v in zip(dates, values):
        deduped[d] = v
    sorted_items = sorted(deduped.items(), key=lambda item: item[0])
    return [item[0] for item in sorted_items], [item[1] for item in sorted_items]




def tail_weekly_series(
    dates: list[date],
    values: list[float | None],
    max_points: int,
) -> tuple[list[date], list[float | None]]:
    """Return the newest date plus the previous dates.

    This is intentionally based on the actual uploaded workbook data, not a
    hardcoded date. If the workbook contains 6/5/2026, the returned window ends
    on 6/5/2026. If a future workbook contains 6/12/2026, it ends there.
    """
    pairs: dict[date, float | None] = {}
    for d, v in zip(dates, values):
        pairs[d] = v
    items = sorted(pairs.items(), key=lambda item: item[0])
    if max_points and max_points > 0:
        items = items[-max_points:]
    return [d for d, _ in items], [v for _, v in items]


def series_to_date_value_map(dates: list[date], values: list[float]) -> dict[date, float]:
    mapped: dict[date, float] = {}
    for d, v in zip(dates, values):
        mapped[d] = v
    return mapped


def newest_date_window(date_lists: list[list[date]], max_points: int) -> list[date]:
    """Build one shared x-axis ending at the newest date found."""
    all_dates = sorted({d for dates in date_lists for d in dates})
    if max_points and max_points > 0:
        all_dates = all_dates[-max_points:]
    return all_dates

def make_tick_indices(total_points: int, max_ticks: int = 12) -> list[int]:
    if total_points <= max_ticks:
        return list(range(total_points))
    step = max(1, math.ceil(total_points / (max_ticks - 1)))
    indices = list(range(0, total_points, step))
    if indices[-1] != total_points - 1:
        indices.append(total_points - 1)
    return sorted(set(indices))


def date_label(d: date) -> str:
    return f"{d.month}/{d.day}/{d.year}"


def render_rebuilt_line_chart(
    output_path: Path,
    dates: list[date],
    series: dict[str, list[float | None]],
    title: str,
    y_label: str,
    percent_axis: bool,
    width_px: int,
    height_px: int,
) -> Path:
    dpi = 160
    fig_w = max(width_px / dpi, 3.0)
    fig_h = max(height_px / dpi, 1.8)
    fig, ax = plt.subplots(figsize=(fig_w, fig_h), dpi=dpi)

    x = list(range(len(dates)))
    for label, values in series.items():
        # Matplotlib will leave a gap for missing values, but the x-axis will
        # still include the newest workbook date.
        ax.plot(x, values, marker="o", linewidth=2.2, markersize=4.2, label=label)

    tick_indices = make_tick_indices(len(dates), max_ticks=11)
    ax.set_xticks(tick_indices)
    ax.set_xticklabels([date_label(dates[i]) for i in tick_indices], rotation=45, ha="right", fontsize=7.5)
    if dates:
        ax.set_xlim(-0.5, len(dates) - 0.5)
    ax.set_title(title, fontsize=12, color="#555555", pad=8)
    ax.set_ylabel(y_label, fontsize=8.5, color="#555555")
    ax.grid(axis="y", alpha=0.28)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.tick_params(axis="y", labelsize=7.5)

    if percent_axis:
        ax.yaxis.set_major_formatter(FuncFormatter(lambda value, _: f"{value:.0%}"))
        all_values = [v for values in series.values() for v in values if v is not None]
        if all_values:
            top = max(max(all_values) * 1.15, 0.10)
            ax.set_ylim(0, min(max(top, 0.10), 1.0))
    else:
        all_values = [v for values in series.values() for v in values if v is not None]
        if all_values:
            top = max(all_values) * 1.18
            ax.set_ylim(0, top if top > 0 else 1)

    if len(series) > 1:
        ax.legend(loc="upper center", bbox_to_anchor=(0.5, -0.25), ncol=min(3, len(series)), fontsize=7, frameon=False)

    fig.tight_layout(pad=1.0)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(output_path, facecolor="white", bbox_inches="tight", pad_inches=0.08)
    plt.close(fig)
    return output_path


def paste_chart(image: Image.Image, chart_path: Path, box: tuple[int, int, int, int]) -> None:
    left, top, right, bottom = box
    width = max(1, right - left)
    height = max(1, bottom - top)
    chart = Image.open(chart_path).convert("RGB").resize((width, height), Image.LANCZOS)
    draw = ImageDraw.Draw(image)
    draw.rectangle(box, fill="white", outline=(225, 225, 225), width=1)
    image.paste(chart, (left, top))


def rebuild_weekly_charts_on_sheet_image(
    image_path: Path,
    output_path: Path,
    prepared_workbook_path: Path,
    sheet_name: str,
    working_dir: Path,
    weekly_window_points: int = 10,
) -> Path:
    """Replace stale Excel-rendered weekly charts with freshly drawn charts.

    This is the important fix for future weekly pulls. Excel/LibreOffice may
    show a chart that visually stops at an older tick label even after formulas
    are extended. This function reads the newest weekly data directly from
    NittyGrittySheet and pastes newly rendered chart images over the two top
    weekly charts, guaranteeing the newest workbook date appears on the x-axis.
    """
    wb = load_workbook(prepared_workbook_path, data_only=True)
    rows = find_weekly_block_rows(wb, sheet_name)
    if rows is None:
        image_path.replace(output_path) if image_path != output_path else None
        return output_path if output_path.exists() else image_path

    date_row, placement_value_row, status_date_row, status_rows = rows
    ws = wb["NittyGrittySheet"]

    placement_dates, placement_values = extract_horizontal_weekly_series(ws, date_row, placement_value_row)
    placement_dates, placement_values = tail_weekly_series(placement_dates, placement_values, weekly_window_points)
    if len(placement_dates) < 2:
        image_path.replace(output_path) if image_path != output_path else None
        return output_path if output_path.exists() else image_path

    status_maps: dict[str, dict[date, float]] = {}
    status_date_lists: list[list[date]] = []
    for label, row in status_rows.items():
        dts, vals = extract_horizontal_weekly_series(ws, status_date_row, row, fallback_date_row=date_row)
        if len(dts) >= 2:
            status_maps[label] = series_to_date_value_map(dts, vals)
            status_date_lists.append(dts)

    status_dates = newest_date_window(status_date_lists, weekly_window_points) if status_date_lists else []
    status_series: dict[str, list[float | None]] = {
        label: [mapped.get(d) for d in status_dates]
        for label, mapped in status_maps.items()
    }

    image = Image.open(image_path).convert("RGB")
    width, height = image.size

    # These boxes match the BYU Marriott placement dashboard screenshot layout.
    left_chart_box = (
        int(width * 0.018),
        int(height * 0.058),
        int(width * 0.490),
        int(height * 0.455),
    )
    right_chart_box = (
        int(width * 0.510),
        int(height * 0.058),
        int(width * 0.982),
        int(height * 0.455),
    )

    chart_dir = working_dir / "rebuilt_weekly_charts"
    placement_chart = chart_dir / f"{sanitize_filename(sheet_name)}_placement.png"
    render_rebuilt_line_chart(
        placement_chart,
        placement_dates,
        {"% Placed": placement_values},
        "Weekly 2026 Placement Trend" if sheet_name == "MSB Full-Time Dashboard" else "Weekly Placement Trend",
        "% of Seeking Students Placed",
        True,
        left_chart_box[2] - left_chart_box[0],
        left_chart_box[3] - left_chart_box[1],
    )
    paste_chart(image, placement_chart, left_chart_box)

    if status_dates and status_series:
        status_chart = chart_dir / f"{sanitize_filename(sheet_name)}_status.png"
        render_rebuilt_line_chart(
            status_chart,
            status_dates,
            status_series,
            "Weekly 2026 Search Status" if sheet_name == "MSB Full-Time Dashboard" else "Weekly Search Status",
            "# of Students",
            False,
            right_chart_box[2] - right_chart_box[0],
            right_chart_box[3] - right_chart_box[1],
        )
        paste_chart(image, status_chart, right_chart_box)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    image.save(output_path, quality=95, optimize=True)
    return output_path

def convert_excel_to_pptx(
    uploaded_excel,
    selected_sheets: list[str],
    uploaded_logo,
    uploaded_intro_pptx,
    title: str,
    fit_each_sheet_to_one_page: bool,
    crop_pages: bool,
    make_logo_transparent: bool,
    show_sheet_name: bool,
    dpi: int,
    sidebar_width: float,
    margins: float,
    auto_extend_latest_date: bool,
    cap_chart_dates_at_today: bool,
    rebuild_weekly_charts: bool,
    weekly_window_points: int,
) -> tuple[bytes, bytes, str]:
    soffice_path = find_soffice()
    if not soffice_path:
        raise RuntimeError(
            "LibreOffice was not found. Install LibreOffice locally, or deploy with packages.txt on Streamlit Cloud."
        )

    with tempfile.TemporaryDirectory() as tmp:
        tmpdir = Path(tmp)
        excel_path = save_upload(uploaded_excel, tmpdir / sanitize_filename(uploaded_excel.name))

        logo_path = None
        if uploaded_logo:
            raw_logo = save_upload(uploaded_logo, tmpdir / sanitize_filename(uploaded_logo.name))
            logo_path = tmpdir / "logo.png"
            if make_logo_transparent:
                make_circular_logo_transparent(raw_logo, logo_path)
            else:
                logo_path = raw_logo
        elif DEFAULT_LOGO.exists():
            logo_path = DEFAULT_LOGO

        intro_template = None
        if uploaded_intro_pptx:
            intro_template = save_upload(uploaded_intro_pptx, tmpdir / sanitize_filename(uploaded_intro_pptx.name))

        prepared_xlsx = tmpdir / "prepared_dashboard.xlsx"
        rendered_sheet_names = prepare_workbook_for_rendering(
            excel_path,
            prepared_xlsx,
            selected_sheets,
            fit_each_sheet_to_one_page,
            margins,
            auto_extend_latest_date,
            cap_chart_dates_at_today,
        )

        pdf_path = convert_to_pdf(prepared_xlsx, tmpdir / "pdf", soffice_path)
        page_images = pdf_to_images(pdf_path, tmpdir / "pages", dpi=dpi)

        if len(page_images) < len(rendered_sheet_names):
            raise RuntimeError(
                f"Expected at least {len(rendered_sheet_names)} rendered pages, but got {len(page_images)}. "
                "Try selecting fewer sheets or disabling one-page fitting."
            )

        final_images = []
        for index, img_path in enumerate(page_images[: len(rendered_sheet_names)], start=1):
            sheet_name = rendered_sheet_names[index - 1]
            working = img_path
            if crop_pages:
                cropped = tmpdir / "cropped" / f"sheet_{index:02d}.png"
                cropped.parent.mkdir(exist_ok=True)
                working = crop_white_space(working, cropped)

            if rebuild_weekly_charts:
                rebuilt = tmpdir / "rebuilt_sheets" / f"sheet_{index:02d}.png"
                rebuilt.parent.mkdir(exist_ok=True)
                try:
                    working = rebuild_weekly_charts_on_sheet_image(
                        working,
                        rebuilt,
                        prepared_xlsx,
                        sheet_name,
                        tmpdir,
                        weekly_window_points=weekly_window_points,
                    )
                except Exception:
                    # Keep the normal Excel screenshot if a workbook does not use
                    # the BYU Marriott NittyGrittySheet dashboard layout.
                    pass

            optimized = tmpdir / "optimized" / f"sheet_{index:02d}.jpg"
            optimized.parent.mkdir(exist_ok=True)
            working = resize_for_ppt(working, optimized)
            final_images.append(working)

        safe_title = sanitize_filename(title.replace(" ", "_"))
        pptx_path = tmpdir / f"{safe_title}.pptx"
        build_powerpoint(
            final_images,
            rendered_sheet_names,
            pptx_path,
            logo_path,
            intro_template,
            tmpdir,
            soffice_path,
            show_sheet_name,
            sidebar_width,
            title,
        )

        zip_path = tmpdir / f"{safe_title}.zip"
        make_zip(pptx_path, zip_path)

        return pptx_path.read_bytes(), zip_path.read_bytes(), pptx_path.name


def get_sheet_names_from_upload(uploaded_excel) -> list[str]:
    with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp:
        tmp.write(uploaded_excel.getbuffer())
        tmp_path = Path(tmp.name)
    try:
        wb = load_workbook(tmp_path, read_only=True)
        return wb.sheetnames
    finally:
        try:
            tmp_path.unlink()
        except OSError:
            pass




def get_newest_date_from_upload(uploaded_excel) -> date | None:
    with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp:
        tmp.write(uploaded_excel.getbuffer())
        tmp_path = Path(tmp.name)
    try:
        wb = load_workbook(tmp_path, read_only=True, data_only=False)
        newest = None
        for ws in wb.worksheets:
            for row in ws.iter_rows():
                for cell in row:
                    parsed = parse_excel_date(cell.value)
                    if parsed and (newest is None or parsed > newest):
                        newest = parsed
        return newest
    finally:
        try:
            tmp_path.unlink()
        except OSError:
            pass


def format_date_for_display(value: date) -> str:
    return f"{value.month}/{value.day}/{value.year}"

def main() -> None:
    st.set_page_config(page_title=APP_TITLE, page_icon="📊", layout="wide")

    st.title("📊 Excel to PowerPoint Dashboard Builder")
    st.caption("Upload an Excel dashboard and turn every worksheet into a clean PowerPoint slide.")

    with st.sidebar:
        st.header("Settings")
        title = st.text_input("Output deck name", value="Placement_Report_Updated")
        fit_one_page = st.checkbox("Fit each worksheet to one landscape page", value=True)
        crop_pages = st.checkbox("Crop extra white space around screenshots", value=True)
        show_sheet_name = st.checkbox("Show sheet name in the right sidebar", value=True)
        make_logo_transparent = st.checkbox("Remove outside background from circular logo", value=True)
        dpi = st.slider("Screenshot resolution", min_value=120, max_value=260, value=190, step=10)
        sidebar_width = st.slider("Right logo sidebar width", min_value=0.8, max_value=1.6, value=1.2, step=0.05)
        margins = st.slider("Excel print margins", min_value=0.05, max_value=0.40, value=0.20, step=0.05)
        auto_extend_latest_date = st.checkbox(
            "Auto-extend charts to newest workbook date",
            value=True,
            help="Uses the newest date found in the uploaded workbook, so future pulls like 5/22, 5/29, 6/5, etc. are included automatically.",
        )
        cap_chart_dates_at_today = st.checkbox(
            "Ignore dates after today",
            value=False,
            help="Leave this off for normal use. Turn it on only if your workbook contains blank future placeholder dates that should not appear yet.",
        )
        rebuild_weekly_charts = st.checkbox(
            "Rebuild weekly charts from NittyGrittySheet",
            value=True,
            help="Recommended. Forces weekly placement and search-status charts to end on the newest date in the uploaded workbook instead of relying on stale Excel chart rendering.",
        )
        weekly_window_points = st.slider(
            "Weekly chart window: newest date plus previous pulls",
            min_value=6,
            max_value=16,
            value=10,
            step=1,
            help="The rebuilt weekly charts will always include the newest workbook date and this many total weekly pulls. Use 10 for roughly the last two months.",
        )

    col1, col2 = st.columns(2)
    with col1:
        excel_upload = st.file_uploader("Upload Excel workbook (.xlsx)", type=["xlsx"])
        logo_upload = st.file_uploader("Optional logo image (.png, .jpg)", type=["png", "jpg", "jpeg"])
    with col2:
        intro_upload = st.file_uploader(
            "Optional PowerPoint template for intro slide (.pptx)",
            type=["pptx"],
            help="If you upload your old deck, the app will use slide 1 as the intro slide.",
        )

    if not excel_upload:
        st.info("Upload an Excel file to begin.")
        return

    try:
        sheet_names = get_sheet_names_from_upload(excel_upload)
    except Exception as exc:
        st.error(f"Could not read the workbook: {exc}")
        return

    newest_workbook_date = get_newest_date_from_upload(excel_upload)
    if newest_workbook_date:
        st.info(f"Newest date detected in this workbook: **{format_date_for_display(newest_workbook_date)}**")

    st.subheader("Sheets to include")
    default_sheets = sheet_names
    selected_sheets = st.multiselect(
        "Choose sheets and order them",
        options=sheet_names,
        default=default_sheets,
    )

    st.write(f"Selected **{len(selected_sheets)}** sheet(s).")

    if st.button("Generate PowerPoint", type="primary", disabled=not selected_sheets):
        try:
            with st.spinner("Rendering Excel sheets and building PowerPoint..."):
                pptx_bytes, zip_bytes, pptx_name = convert_excel_to_pptx(
                    uploaded_excel=excel_upload,
                    selected_sheets=selected_sheets,
                    uploaded_logo=logo_upload,
                    uploaded_intro_pptx=intro_upload,
                    title=title,
                    fit_each_sheet_to_one_page=fit_one_page,
                    crop_pages=crop_pages,
                    make_logo_transparent=make_logo_transparent,
                    show_sheet_name=show_sheet_name,
                    dpi=dpi,
                    sidebar_width=sidebar_width,
                    margins=margins,
                    auto_extend_latest_date=auto_extend_latest_date,
                    cap_chart_dates_at_today=cap_chart_dates_at_today,
                    rebuild_weekly_charts=rebuild_weekly_charts,
                    weekly_window_points=weekly_window_points,
                )

            st.success("PowerPoint created successfully.")
            st.download_button(
                "Download PowerPoint",
                data=pptx_bytes,
                file_name=pptx_name,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
            st.download_button(
                "Download ZIP",
                data=zip_bytes,
                file_name=pptx_name.replace(".pptx", ".zip"),
                mime="application/zip",
            )
        except Exception as exc:
            st.error(str(exc))
            st.caption(
                "Common fix: make sure LibreOffice is installed locally, or include packages.txt when deploying to Streamlit Cloud."
            )


if __name__ == "__main__":
    main()
